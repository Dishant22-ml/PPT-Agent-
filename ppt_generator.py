"""
XML to PPTX Converter
Converts modified XML back to PowerPoint presentation
Reverses the extraction process from slide_extractor.py
"""

import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path
import sys


class XMLToPPTXConverter:
    """Convert modified XML back to PPTX"""
    
    def __init__(self):
        self.slide_width_emu = 9144000
        self.slide_height_emu = 6858000
    
    def convert(self, original_pptx: str, modified_xml: str, output_pptx: str):
        """
        Convert modified XML back to PPTX
        
        Args:
            original_pptx: Path to original PPTX (used as template)
            modified_xml: Path to modified XML with changes
            output_pptx: Path for output PPTX file
        """
        print(f"\n{'='*60}")
        print("XML TO PPTX CONVERSION")
        print(f"{'='*60}\n")
        
        # Load original presentation as template
        print(f"📂 Loading template: {original_pptx}")
        prs = Presentation(original_pptx)
        
        # Parse modified XML
        print(f"📄 Parsing modified XML: {modified_xml}")
        tree = ET.parse(modified_xml)
        root = tree.getroot()
        
        # Extract slide dimensions from XML if available
        dimensions = root.find('.//slide_dimensions')
        if dimensions is not None:
            width_elem = dimensions.find('width')
            height_elem = dimensions.find('height')
            if width_elem is not None and height_elem is not None:
                self.slide_width_emu = int(width_elem.text)
                self.slide_height_emu = int(height_elem.text)
        
        # Get slide element (for single slide analysis)
        xml_slide = root.find('.//slide')
        if xml_slide is None:
            print("❌ No slide found in XML")
            return
        
        slide_index = int(xml_slide.get('index', 1)) - 1  # Convert to 0-based
        
        if slide_index >= len(prs.slides):
            print(f"⚠️  Warning: Slide index {slide_index} out of range")
            return
        
        print(f"📊 Processing Slide {slide_index + 1}...\n")
        slide = prs.slides[slide_index]
        
        # Apply modifications
        self._apply_slide_modifications(slide, xml_slide)
        
        # Save modified presentation
        print(f"\n💾 Saving to: {output_pptx}")
        prs.save(output_pptx)
        
        print(f"\n{'='*60}")
        print(f"✅ SUCCESS! PPTX saved to: {output_pptx}")
        print(f"{'='*60}\n")
    
    def _apply_slide_modifications(self, slide, xml_slide: ET.Element):
        """Apply modifications from XML to PowerPoint slide"""
        
        elements = xml_slide.find('.//elements')
        if elements is None:
            print("  ⚠️  No elements found in XML")
            return
        
        # Track which shapes we've updated
        updated_shapes = set()
        
        # Process each modified element
        for xml_elem in elements.findall('.//element'):
            elem_id = xml_elem.get('id')
            elem_type = xml_elem.get('type')
            
            # Find corresponding shape in slide
            shape = self._find_shape_by_id(slide, elem_id)
            
            if shape is None:
                print(f"  ⚠️  Shape {elem_id} not found, skipping")
                continue
            
            # Apply modifications based on type
            if elem_type in ['textbox', 'shape']:
                self._update_text_shape(shape, xml_elem)
                updated_shapes.add(elem_id)
                print(f"  ✓ Updated text in shape: {elem_id}")
                
            elif elem_type == 'picture':
                # Image updates would require new image file handling
                self._update_picture(shape, xml_elem)
                updated_shapes.add(elem_id)
                print(f"  ✓ Updated picture: {elem_id}")
        
        print(f"\n  Total shapes updated: {len(updated_shapes)}")
    
    def _find_shape_by_id(self, slide, elem_id: str):
        """Find shape in slide by element ID"""
        # Try to match by shape ID or name
        for shape in slide.shapes:
            # Match by ID
            if str(shape.shape_id) == elem_id:
                return shape
            
            # Match by name
            if hasattr(shape, 'name') and shape.name == elem_id:
                return shape
            
            # Try matching last part of ID if it's numeric
            try:
                if str(shape.shape_id) == elem_id.split('_')[-1]:
                    return shape
            except:
                pass
        
        return None
    
    def _update_text_shape(self, shape, xml_elem: ET.Element):
        """Update text content in a shape"""
        # Check if shape has text frame
        if not hasattr(shape, 'text_frame'):
            return
        
        text_frame = shape.text_frame
        
        # Get text body from XML
        text_body = xml_elem.find('.//text_body')
        if text_body is None:
            return
        
        # Clear existing text
        text_frame.clear()
        
        # Process each paragraph
        paragraphs = text_body.findall('.//paragraph')
        for para_elem in paragraphs:
            # Add paragraph
            if len(text_frame.paragraphs) == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # Set paragraph alignment
            alignment = para_elem.find('.//alignment')
            if alignment is not None and alignment.text:
                align_map = {
                    'left': PP_ALIGN.LEFT,
                    'ctr': PP_ALIGN.CENTER,
                    'center': PP_ALIGN.CENTER,
                    'r': PP_ALIGN.RIGHT,
                    'right': PP_ALIGN.RIGHT,
                    'just': PP_ALIGN.JUSTIFY
                }
                p.alignment = align_map.get(alignment.text.lower(), PP_ALIGN.LEFT)
            
            # Set paragraph level (for bullets)
            level = para_elem.get('level', '0')
            p.level = int(level)
            
            # Process text runs
            text_runs = para_elem.findall('.//text_run')
            for run_elem in text_runs:
                text_elem = run_elem.find('.//text')
                if text_elem is None or text_elem.text is None:
                    continue
                
                text = text_elem.text
                
                # Add run
                run = p.add_run()
                run.text = text
                
                # Apply formatting
                font_elem = run_elem.find('.//font')
                if font_elem is not None:
                    # Font family
                    family = font_elem.get('family')
                    if family:
                        run.font.name = family
                    
                    # Font size
                    size = font_elem.get('size')
                    if size:
                        try:
                            run.font.size = Pt(float(size))
                        except:
                            pass
                    
                    # Bold, italic, underline
                    bold = font_elem.get('bold', 'false').lower() == 'true'
                    italic = font_elem.get('italic', 'false').lower() == 'true'
                    underline = font_elem.get('underline', 'false').lower() == 'true'
                    
                    run.font.bold = bold
                    run.font.italic = italic
                    run.font.underline = underline
                
                # Apply color
                color_elem = run_elem.find('.//color')
                if color_elem is not None:
                    hex_color = color_elem.get('hex')
                    if hex_color and hex_color.startswith('#'):
                        try:
                            rgb = self._hex_to_rgb(hex_color)
                            from pptx.util import RGBColor
                            run.font.color.rgb = RGBColor(*rgb)
                        except:
                            pass
    
    def _update_picture(self, shape, xml_elem: ET.Element):
        """Update picture properties (metadata only, no image replacement)"""
        # In a full implementation, you would:
        # 1. Check for new_image_desc attribute (from LLM modifications)
        # 2. Generate or fetch new image
        # 3. Replace the image in the shape
        
        # For now, just update alt text if available
        acc = xml_elem.find('.//accessibility/alt_text')
        if acc is not None and acc.text:
            if hasattr(shape, 'element'):
                # Update alt text in underlying XML
                try:
                    shape.element.cNvPr.set('descr', acc.text)
                except:
                    pass
    
    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def main():
    """Main entry point"""
    if len(sys.argv) < 4:
        print("Usage: python xml_to_pptx.py <original.pptx> <modified.xml> <output.pptx>")
        print("\nExample:")
        print("  python xml_to_pptx.py template.pptx presentation_modified.xml output.pptx")
        return
    
    original_pptx = sys.argv[1]
    modified_xml = sys.argv[2]
    output_pptx = sys.argv[3]
    
    # Check files exist
    if not Path(original_pptx).exists():
        print(f"❌ Error: Original PPTX not found: {original_pptx}")
        return
    
    if not Path(modified_xml).exists():
        print(f"❌ Error: Modified XML not found: {modified_xml}")
        return
    
    # Convert
    converter = XMLToPPTXConverter()
    
    try:
        converter.convert(original_pptx, modified_xml, output_pptx)
        print("\n🎉 Conversion complete!")
        
    except Exception as e:
        print(f"\n❌ Error during conversion: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main()